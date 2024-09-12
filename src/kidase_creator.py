import os
import re
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Emu, Pt
from tqdm import tqdm


class KidaseCreator:
    '''
    This class is responsible for aggregating the kidase data and components 
    to create the slides for the Liturgy, readings, and hymns.
    '''
    def __init__(self, data_path, languages):
        self.order_of_the_liturgy = pd.read_excel(os.path.join(data_path, 'ሥርዓተ_ቅዳሴ.xlsx'))
        self.languages = languages
        self.slide_settings = {
            'background_color': RGBColor(0, 0, 0),
            'font': 'Arial',
            'font_size': Pt(24),
            'font_color': RGBColor(255, 255, 255),  
            'border_color': RGBColor(255, 255, 255),
            'border_thickness': Pt(1)  
        }
        self.keyword_mapping = {
            '፠ Priest:': RGBColor(255, 0, 0),  # Red
            '፠ ካህን፤': RGBColor(255, 0, 0),    # Red
            '፠ Asst. Priest:': RGBColor(255, 0, 0),  # Red
            '፠ ካህን ንፍቅ፤': RGBColor(255, 0, 0),  # Red
            '፠ Deacon:': RGBColor(0, 255, 0),  # Green
            '፠ ዲያቆን፤': RGBColor(0, 255, 0),  # Green
            '፠ Asst. Deacon:': RGBColor(0, 255, 0),  # Green
            '፠ ዲያቆን ንፍቅ፤': RGBColor(0, 255, 0),  # Green
            '፠ People:': RGBColor(255, 215, 0),  # Gold
            '፠ ሕዝብ፤': RGBColor(255, 215, 0),  # Gold
            '፠ ALL:': RGBColor(255, 215, 0),     # Gold
            '፠ ኵሎሙ፤': RGBColor(255, 215, 0)  # Gold
        }
        self.process_data()

    def process_data(self):
        '''
        Checks the order of the liturgy for the given languages.
        '''
        for lang in self.languages:
            if lang not in self.order_of_the_liturgy.columns:
                raise ValueError(f"Language '{lang}' not found in the data.")
            self.order_of_the_liturgy[lang] = self.order_of_the_liturgy[lang].fillna('')

    def format_slide(self, slide, slide_settings):
        '''
        Prepare the slide for the given languages.
        '''
        left, top, width, height = slide_settings['text_box_position']
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.word_wrap = True

        text = slide_settings['text']

        # Create a single paragraph for the entire text
        p = text_frame.add_paragraph()
        keywords = sorted(self.keyword_mapping.keys())
        pattern = '|'.join(map(re.escape, keywords))

        split_string = re.split(f'({pattern})', text)
        split_string = [s for s in split_string if s]

        # Go through each keyword and format it accordingly in the text
        for string in split_string:
            if string in self.keyword_mapping:
                run = p.add_run()
                run.text = string
                run.font.color.rgb = self.keyword_mapping[string]
            else:
                run = p.add_run()
                run.text = string
                run.font.color.rgb = slide_settings['font_color']
            run.font.name = slide_settings['font']
            run.font.size = slide_settings['font_size']
        
        # Apply border settings for the text box
        textbox.line.color.rgb = slide_settings['border_color']
        textbox.line.width = slide_settings['border_thickness']


    def create_presentation(self):
        # Create a presentation object
        prs = Presentation()
        prs.slide_height = Emu(5.5 * 914400) # 5.5 inches to EMU
        slide_layout = prs.slide_layouts[6]  # Blank layout
        for i in tqdm(range(len(self.order_of_the_liturgy)), desc="Creating slides"):
            slide = prs.slides.add_slide(slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = self.slide_settings['background_color']  # Black color
            if len(self.languages) in [1, 2, 3]:
                # Use columns for up to 3 languages
                width = prs.slide_width // len(self.languages)
                for j, lang in enumerate(self.languages):
                    self.slide_settings['text_box_position'] = (j * width, 0, width, prs.slide_height)
                    self.slide_settings['text'] = self.order_of_the_liturgy[lang][i]
                    self.format_slide(slide, self.slide_settings)
            elif len(self.languages) == 4:
                # Use 2x2 grid for 4 languages
                width = prs.slide_width // 2
                height = prs.slide_height // 2
                for j, lang in enumerate(self.languages):
                    self.slide_settings['text_box_position'] = ((j % 2) * width, (j // 2) * height, width, height)
                    self.slide_settings['text'] = self.order_of_the_liturgy[lang][i]
                    self.format_slide(slide, self.slide_settings)
            else:
                raise ValueError("Only up to 4 languages are supported.")
        # Save the presentation
        prs.save('test.pptx')


# Example usage with Geez, Tigrinya, and English languages
if __name__ == '__main__':
    kidase_creator = KidaseCreator('/mnt/c/Users/samg/Documents/Gits/KidaseCreator/data', ['ግእዝ', 'ትግርኛ', 'english'])
    kidase_creator.create_presentation()
